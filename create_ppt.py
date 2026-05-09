"""Generate ADMET Risk Profiler presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_BG = RGBColor(0x06, 0x0B, 0x17)
SURFACE = RGBColor(0x0D, 0x15, 0x26)
SURFACE2 = RGBColor(0x11, 0x1D, 0x35)
PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0xA7, 0x8B, 0xFA)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
TEXT2 = RGBColor(0x94, 0xA3, 0xB8)
TEXT3 = RGBColor(0x64, 0x74, 0x8B)
PASS = RGBColor(0x00, 0xE5, 0xA0)
WARN = RGBColor(0xFF, 0xB5, 0x47)
FAIL = RGBColor(0xFF, 0x4D, 0x6D)
TEAL = RGBColor(0x38, 0xBD, 0xF8)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = 0
    return shape


def add_text_box(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════
# SLIDE 1 — Title / Group Details
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative circles
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x4C, 0x1D, 0x95))
add_circle(slide, Inches(10), Inches(4.5), Inches(4), RGBColor(0x1E, 0x3A, 0x5F))
add_circle(slide, Inches(5), Inches(5), Inches(3), RGBColor(0x0F, 0x76, 0x6E))

# Badge
add_rect(slide, Inches(4.5), Inches(0.8), Inches(4.3), Inches(0.45), SURFACE2)
add_text_box(slide, Inches(4.5), Inches(0.82), Inches(4.3), Inches(0.45),
             "PHARMACOLOGY TOOL", size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Title
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
             "ADMET Risk Profiler", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(2), Inches(2.6), Inches(9.3), Inches(0.7),
             "Advanced Drug-Likeness Evaluation Based on Lipinski's Rule of Five", size=18, color=TEXT2, align=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.5), Inches(3.3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY
line.line.fill.background()

# Group heading
add_text_box(slide, Inches(3), Inches(3.8), Inches(7.3), Inches(0.55),
             "PROJECT TEAM", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Member cards
members = ["Ayan Mondal", "Debopriya Saha", "Debraj De", "Krish"]
card_w = Inches(2.8)
card_h = Inches(1.1)
gap = Inches(0.3)
total_w = card_w * 4 + gap * 3
start_x = (prs.slide_width - total_w) // 2

for i, name in enumerate(members):
    x = start_x + i * (card_w + gap)
    y = Inches(4.5)
    card = add_rect(slide, x, y, card_w, card_h, SURFACE2)
    # Icon circle
    add_circle(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.7), PRIMARY)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(0.7), Inches(0.5),
                 str(i + 1), size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide, x + Inches(1), y + Inches(0.25), Inches(1.6), Inches(0.6),
                 name, size=16, color=WHITE, bold=True)

# Footer
add_text_box(slide, Inches(2), Inches(6.5), Inches(9.3), Inches(0.5),
             "B.Pharm Project  |  2026", size=14, color=TEXT3, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 2 — What is ADMET?
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(11), Inches(-1), Inches(3), RGBColor(0x4C, 0x1D, 0x95))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "What is ADMET?", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "The five pillars of pharmacokinetic drug evaluation", size=16, color=TEXT2)

items = [
    ("A", "Absorption", "GI absorption, water solubility, bioavailability score, MW, and PSA determine oral uptake", PASS),
    ("D", "Distribution", "LogP, PSA, and BBB permeability govern tissue and CNS distribution", TEAL),
    ("M", "Metabolism", "Rotatable bonds, bioavailability score & water solubility predict metabolic stability", ACCENT),
    ("E", "Excretion", "How the kidneys and other organs clear the drug from the body", WARN),
    ("T", "Toxicity", "Structural flags for potential toxic or adverse effects", FAIL),
]

for i, (letter, title, desc, color) in enumerate(items):
    y = Inches(2.2 + i * 1.0)
    add_circle(slide, Inches(1), y, Inches(0.65), color)
    add_text_box(slide, Inches(1), y + Inches(0.1), Inches(0.65), Inches(0.5),
                 letter, size=22, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), y + Inches(0.05), Inches(4), Inches(0.45),
                 title, size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2), y + Inches(0.45), Inches(9), Inches(0.45),
                 desc, size=14, color=TEXT2)


# ════════════════════════════════════════════
# SLIDE 3 — Lipinski's Rule of Five
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(-1.5), Inches(5), Inches(4), RGBColor(0x1E, 0x3A, 0x5F))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Lipinski's Rule of Five", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "The gold standard for predicting oral bioavailability of drug candidates", size=16, color=TEXT2)

rules = [
    ("Molecular Weight", "≤ 500 Da", "Sum of atomic masses of all atoms in the molecule"),
    ("LogP (Lipophilicity)", "≤ 5", "Octanol-water partition coefficient; ideal range 1–3"),
    ("H-Bond Donors", "≤ 5", "Count of –OH and –NH groups in the molecule"),
    ("H-Bond Acceptors", "≤ 10", "Count of nitrogen and oxygen atoms"),
]

for i, (param, limit, desc) in enumerate(rules):
    y = Inches(2.3 + i * 1.15)
    card = add_rect(slide, Inches(1), y, Inches(11.3), Inches(0.95), SURFACE2)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(4), Inches(0.4),
                 param, size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.5), Inches(7), Inches(0.35),
                 desc, size=13, color=TEXT2)
    # Limit badge
    badge = add_rect(slide, Inches(10), y + Inches(0.2), Inches(2), Inches(0.5), PRIMARY)
    add_text_box(slide, Inches(10), y + Inches(0.25), Inches(2), Inches(0.45),
                 limit, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
             "⚠ Two or more violations suggest poor oral bioavailability — compound is high risk.",
             size=14, color=WARN, bold=False)


# ════════════════════════════════════════════
# SLIDE 4 — Technology Stack
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(10.5), Inches(5), Inches(3.5), RGBColor(0x0F, 0x76, 0x6E))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Technology Stack", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Full-stack web application with serverless deployment", size=16, color=TEXT2)

tech = [
    ("Frontend", "HTML5, CSS3, JavaScript, Chart.js", "Modern dark UI with glassmorphism, radar charts, and micro-animations", TEAL),
    ("Backend", "Python, Flask, Flask-CORS", "RESTful API for ADMET evaluation with scoring engine", PRIMARY),
    ("Deployment", "Vercel (Serverless)", "Auto-detects local vs production environment for API routing", ACCENT),
    ("Algorithm", "Lipinski + Veber Rules", "Multi-parameter scoring across all five ADMET dimensions", PASS),
]

for i, (title, stack, desc, color) in enumerate(tech):
    y = Inches(2.2 + i * 1.2)
    card = add_rect(slide, Inches(1), y, Inches(11.3), Inches(1.0), SURFACE2)
    # Color accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y, Pt(6), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, Inches(1.4), y + Inches(0.1), Inches(3), Inches(0.4),
                 title, size=20, color=color, bold=True)
    add_text_box(slide, Inches(4.5), y + Inches(0.12), Inches(7), Inches(0.35),
                 stack, size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.55), Inches(10), Inches(0.35),
                 desc, size=13, color=TEXT2)


# ════════════════════════════════════════════
# SLIDE 5 — How It Works
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(-1), Inches(-1), Inches(3), RGBColor(0x4C, 0x1D, 0x95))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "How It Works", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "User workflow from input to analysis results", size=16, color=TEXT2)

steps = [
    ("1", "Input Properties", "Enter MW, LogP, H-bond donors/acceptors, rotatable bonds, PSA,\nGI absorption, water solubility, bioavailability score & BBB permeability"),
    ("2", "Run Analysis", "Flask backend evaluates compound against Lipinski's Rule\nof Five and extended ADMET criteria"),
    ("3", "View Results", "Dashboard displays Pass/Warning/Fail for each ADMET\ncategory with scores, radar chart, and actionable tips"),
    ("4", "Track History", "Compare multiple compounds side-by-side in the\nanalysis history panel"),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.2 + i * 1.2)
    # Number circle
    add_circle(slide, Inches(1.2), y + Inches(0.1), Inches(0.7), PRIMARY)
    add_text_box(slide, Inches(1.2), y + Inches(0.2), Inches(0.7), Inches(0.5),
                 num, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Connector line
    if i < 3:
        conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.53), y + Inches(0.85), Pt(3), Inches(0.5))
        conn.fill.solid()
        conn.fill.fore_color.rgb = PRIMARY
        conn.line.fill.background()
    add_text_box(slide, Inches(2.3), y + Inches(0.05), Inches(4), Inches(0.4),
                 title, size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.3), y + Inches(0.45), Inches(9), Inches(0.6),
                 desc, size=13, color=TEXT2)


# ════════════════════════════════════════════
# SLIDE 6 — Scoring System
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(11), Inches(5.5), Inches(3), RGBColor(0x1E, 0x3A, 0x5F))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Scoring System", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Each ADMET dimension is scored 0–100 and classified into three tiers", size=16, color=TEXT2)

tiers = [
    ("PASS  (80–100)", "Compound meets or exceeds safety thresholds", PASS, "✅"),
    ("WARNING  (55–79)", "Borderline — compound may have pharmacokinetic issues", WARN, "⚠️"),
    ("FAIL  (0–54)", "High risk — compound violates critical safety rules", FAIL, "❌"),
]

for i, (tier, desc, color, icon) in enumerate(tiers):
    y = Inches(2.3 + i * 1.3)
    card = add_rect(slide, Inches(1.5), y, Inches(10.3), Inches(1.05), SURFACE2)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y, Pt(6), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, Inches(2), y + Inches(0.12), Inches(5), Inches(0.45),
                 f"{icon}  {tier}", size=20, color=color, bold=True)
    add_text_box(slide, Inches(2), y + Inches(0.55), Inches(9), Inches(0.4),
                 desc, size=14, color=TEXT2)

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.6),
             "Overall Drug-Likeness Score = Average of all 5 ADMET dimension scores\nLipinski Violations counted separately for quick risk flagging",
             size=13, color=TEXT3)


# ════════════════════════════════════════════
# SLIDE 7 — Input Parameters
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Input Parameters", size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
             "Ten molecular descriptors power the ADMET analysis engine", size=16, color=TEXT2)

params = [
    ("Molecular Weight", "Da", "≤ 500"),
    ("LogP", "—", "≤ 5"),
    ("H-Bond Donors", "count", "≤ 5"),
    ("H-Bond Acceptors", "count", "≤ 10"),
    ("Rotatable Bonds", "count", "≤ 10"),
    ("Polar Surface Area", "Å²", "≤ 140"),
    ("GI Absorption", "—", "High"),
    ("Water Solubility", "—", "High"),
    ("Bioavailability Score", "0–1", "≥ 0.55"),
    ("BBB Permeability", "—", "Yes/No"),
]

# Two columns of 5
for i, (name, unit, limit) in enumerate(params):
    col = i // 5
    row = i % 5
    x = Inches(0.8 + col * 6.3)
    y = Inches(2.0 + row * 1.05)
    card = add_rect(slide, x, y, Inches(5.8), Inches(0.88), SURFACE2)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.08), Inches(3.5), Inches(0.38),
                 name, size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.48), Inches(2), Inches(0.33),
                 f"Unit: {unit}", size=11, color=TEXT3)
    badge = add_rect(slide, x + Inches(4.1), y + Inches(0.2), Inches(1.4), Inches(0.45), PRIMARY)
    add_text_box(slide, x + Inches(4.1), y + Inches(0.23), Inches(1.4), Inches(0.4),
                 limit, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# SLIDE 8 — Key Features
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(10), Inches(-1), Inches(3.5), RGBColor(0x4C, 0x1D, 0x95))

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Key Features", size=36, color=WHITE, bold=True)

features = [
    ("Real-time Visual Feedback", "Live range bars and animated score indicators as you type"),
    ("Radar Chart Visualization", "Interactive Chart.js radar chart for ADMET profile at a glance"),
    ("Detailed Risk Analysis", "Expandable cards with explanations, tips, and rule references"),
    ("Analysis History", "Track and compare up to 6 compound evaluations per session"),
    ("Responsive Dark UI", "Glassmorphism design with animated orbs and micro-animations"),
    ("Serverless Deployment", "Vercel-ready with auto-detecting local/production API routing"),
]

for i, (title, desc) in enumerate(features):
    col = i // 3
    row = i % 3
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 1.7)
    card = add_rect(slide, x, y, Inches(5.8), Inches(1.4), SURFACE2)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.2), Inches(5), Inches(0.45),
                 title, size=18, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.4), y + Inches(0.7), Inches(5), Inches(0.5),
                 desc, size=14, color=TEXT2)


# ════════════════════════════════════════════
# SLIDE 9 — Test Cases
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Test Cases", size=36, color=WHITE, bold=True)

# Good drug
card = add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(5.2), SURFACE2)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.8), Pt(5))
bar.fill.solid(); bar.fill.fore_color.rgb = PASS; bar.line.fill.background()
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(5), Inches(0.5),
             "✅ Aspirin — Perfect Drug", size=20, color=PASS, bold=True)
good_vals = "Molecular Weight:  180 Da\nLogP:  1.2\nH-Bond Donors:  1\nH-Bond Acceptors:  4\nRotatable Bonds:  3\nPolar Surface Area:  63.6 Å²\nGI Absorption:  High\nWater Solubility:  High\nBioavailability Score:  0.85\nBBB Permeability:  No"
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(5), Inches(3.5),
             good_vals, size=13, color=TEXT2)
add_text_box(slide, Inches(1.2), Inches(6.0), Inches(5), Inches(0.4),
             "Expected: All PASS  |  Score: ~100", size=13, color=PASS, bold=True)

# Bad drug
card = add_rect(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), SURFACE2)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.8), Pt(5))
bar.fill.solid(); bar.fill.fore_color.rgb = FAIL; bar.line.fill.background()
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
             "❌ BadDrug-X — High Risk", size=20, color=FAIL, bold=True)
bad_vals = "Molecular Weight:  750 Da\nLogP:  6.5\nH-Bond Donors:  8\nH-Bond Acceptors:  14\nRotatable Bonds:  12\nPolar Surface Area:  180 Å²\nGI Absorption:  Low\nWater Solubility:  Poor\nBioavailability Score:  0.20\nBBB Permeability:  Yes"
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(3.5),
             bad_vals, size=13, color=TEXT2)
add_text_box(slide, Inches(7.2), Inches(6.0), Inches(5), Inches(0.4),
             "Expected: Multiple FAIL  |  Score: ~15", size=13, color=FAIL, bold=True)


# ════════════════════════════════════════════
# SLIDE 10 — Thank You
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_circle(slide, Inches(0), Inches(3), Inches(5), RGBColor(0x4C, 0x1D, 0x95))
add_circle(slide, Inches(9), Inches(1), Inches(4), RGBColor(0x1E, 0x3A, 0x5F))
add_circle(slide, Inches(5), Inches(5), Inches(3), RGBColor(0x0F, 0x76, 0x6E))

add_text_box(slide, Inches(1.5), Inches(2), Inches(10.3), Inches(1.5),
             "Thank You!", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.7),
             "ADMET Risk Profiler — For Educational Use Only", size=18, color=ACCENT, align=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Pt(2))
line.fill.solid(); line.fill.fore_color.rgb = PRIMARY; line.line.fill.background()

add_text_box(slide, Inches(2), Inches(4.6), Inches(9.3), Inches(1.2),
             "Ayan Mondal  •  Debopriya Saha  •  Debraj De  •  Krish",
             size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
             "Questions?", size=16, color=TEXT2, align=PP_ALIGN.CENTER)


# ── Save ──
output = r"c:\Users\pitay\Desktop\ADMET\ADMET_Risk_Profiler.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
